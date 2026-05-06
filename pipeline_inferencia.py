"""
pipeline_inferencia.py
Pipeline de inferencia modal — TFM Simulador de Movilidad Urbana
Genera pipeline_inferencia.pptx  (28 × 18 cm, fondo blanco, Calibri)
Ejecutar: python3 pipeline_inferencia.py
"""
from pptx import Presentation
from pptx.util import Cm, Pt
from pptx.dml.color import RGBColor
from pptx.enum.text import PP_ALIGN, MSO_ANCHOR
from pptx.oxml.ns import qn
import lxml.etree as etree

# ── Slide ─────────────────────────────────────────────────────────────────────
SLIDE_W = Cm(28)
SLIDE_H = Cm(18)

prs = Presentation()
prs.slide_width  = SLIDE_W
prs.slide_height = SLIDE_H
slide = prs.slides.add_slide(prs.slide_layouts[6])  # blank

# ── Paleta ────────────────────────────────────────────────────────────────────
C_CYAN   = RGBColor(0x00, 0xBC, 0xD4)
C_ORANGE = RGBColor(0xFF, 0x8F, 0x00)
C_BGRAY  = RGBColor(0x54, 0x6E, 0x7A)
C_DKGR   = RGBColor(0x2E, 0x7D, 0x32)
C_DKBL   = RGBColor(0x15, 0x65, 0xC0)
C_GREEN  = RGBColor(0x43, 0xA0, 0x47)
C_RED    = RGBColor(0xE5, 0x39, 0x35)
C_BLACK  = RGBColor(0x00, 0x00, 0x00)
C_DGRAY  = RGBColor(0x44, 0x44, 0x44)
C_WHITE  = RGBColor(0xFF, 0xFF, 0xFF)

def tint(rgb, f=0.92):
    return RGBColor(
        int(rgb[0] + (255 - rgb[0]) * f),
        int(rgb[1] + (255 - rgb[1]) * f),
        int(rgb[2] + (255 - rgb[2]) * f),
    )

# ── Layout ────────────────────────────────────────────────────────────────────
CX     = 14.0   # centro horizontal (cm)
MAIN_W = 13.0   # ancho bloques principales
MAIN_H = 1.42   # alto bloques principales
AW     = Pt(2)  # grosor de flecha

# Posiciones Y (cm) — layout distribuido en 18 cm
Y_B1    = 0.82                       # Input
B1_BOT  = Y_B1 + MAIN_H             # = 2.24

Y_HUB1  = B1_BOT + 0.52             # = 2.76  nodo fan-out asyncio
Y_PAR1  = Y_HUB1 + 0.42             # = 3.18  top 4 cajas paralelas
H_PAR1  = 1.68
PAR1_BOT = Y_PAR1 + H_PAR1          # = 4.86

Y_HUB2  = PAR1_BOT + 0.40           # = 5.26  nodo fan-in
Y_B3    = Y_HUB2 + 0.38             # = 5.64  build_route_features
B3_BOT  = Y_B3 + MAIN_H             # = 7.06

Y_B4    = B3_BOT + 0.62             # = 7.68  build_feature_frame
B4_BOT  = Y_B4 + MAIN_H             # = 9.10

Y_B5    = B4_BOT + 0.62             # = 9.72  StandardScaler
B5_BOT  = Y_B5 + MAIN_H             # = 11.14

Y_HUB3  = B5_BOT + 0.50             # = 11.64 nodo fan-out modelos
Y_PAR2  = Y_HUB3 + 0.42             # = 12.06 top 3 cajas modelos
H_PAR2  = 1.70
MOD_BOT = Y_PAR2 + H_PAR2           # = 13.76

Y_HUB4  = MOD_BOT + 0.42            # = 14.18 nodo fan-in
Y_B7    = Y_HUB4 + 0.38             # = 14.56 Respuesta JSON
# B7_BOT = 15.98  — margen inferior ~2 cm en slide de 18 cm ✓

# ── Helpers ───────────────────────────────────────────────────────────────────

def set_radius(shape, val=18000):
    sp = shape._element
    spPr = sp.find(qn('p:spPr'))
    if spPr is None:
        return
    pg = spPr.find(qn('a:prstGeom'))
    if pg is None:
        return
    al = pg.find(qn('a:avLst'))
    if al is None:
        al = etree.SubElement(pg, qn('a:avLst'))
    for g in al.findall(qn('a:gd')):
        al.remove(g)
    gd = etree.SubElement(al, qn('a:gd'))
    gd.set('name', 'adj')
    gd.set('fmla', f'val {val}')


def box(x, y, w, h, color, title, subtitle=None, ts=13, ss=9.5, bg=None):
    shape = slide.shapes.add_shape(5, Cm(x), Cm(y), Cm(w), Cm(h))
    set_radius(shape, 16000)
    shape.fill.solid()
    shape.fill.fore_color.rgb = bg if bg else tint(color, 0.93)
    shape.line.color.rgb = color
    shape.line.width = Pt(2.5)
    tf = shape.text_frame
    tf.word_wrap = True
    tf.margin_left = tf.margin_right = Cm(0.30)
    tf.margin_top = tf.margin_bottom = Cm(0.08)
    tf.vertical_anchor = MSO_ANCHOR.MIDDLE
    p1 = tf.paragraphs[0]
    p1.alignment = PP_ALIGN.CENTER
    r1 = p1.add_run()
    r1.text = title
    r1.font.name = 'Calibri'
    r1.font.size = Pt(ts)
    r1.font.bold = True
    r1.font.color.rgb = C_BLACK
    if subtitle:
        p2 = tf.add_paragraph()
        p2.alignment = PP_ALIGN.CENTER
        r2 = p2.add_run()
        r2.text = subtitle
        r2.font.name = 'Calibri'
        r2.font.size = Pt(ss)
        r2.font.bold = False
        r2.font.color.rgb = C_DGRAY
    return shape


def dashed_box(x, y, w, h, color, lbl):
    """Recuadro discontinuo de agrupación con label centrada encima."""
    shape = slide.shapes.add_shape(1, Cm(x), Cm(y), Cm(w), Cm(h))  # RECTANGLE
    shape.fill.background()
    shape.line.color.rgb = color
    shape.line.width = Pt(1.5)
    # Dash via XML
    spPr = shape._element.find(qn('p:spPr'))
    if spPr is not None:
        lnElem = spPr.find(qn('a:ln'))
        if lnElem is None:
            lnElem = etree.SubElement(spPr, qn('a:ln'))
        pd = lnElem.find(qn('a:prstDash'))
        if pd is None:
            pd = etree.SubElement(lnElem, qn('a:prstDash'))
        pd.set('val', 'dashDot')
    # Label centrada encima del recuadro
    lbl_w = 10.0
    txBox = slide.shapes.add_textbox(
        Cm(x + w/2 - lbl_w/2), Cm(y - 0.34), Cm(lbl_w), Cm(0.38))
    tf = txBox.text_frame
    p = tf.paragraphs[0]
    p.alignment = PP_ALIGN.CENTER
    r = p.add_run()
    r.text = lbl
    r.font.name = 'Calibri'
    r.font.size = Pt(10.5)
    r.font.bold = True
    r.font.italic = True
    r.font.color.rgb = color


def line(x1, y1, x2, y2, color=None, w=None, arrow=False):
    color = color or C_BLACK
    w = w or AW
    conn = slide.shapes.add_connector(1, Cm(x1), Cm(y1), Cm(x2), Cm(y2))
    conn.line.color.rgb = color
    conn.line.width = w
    if arrow:
        ln = conn.line._get_or_add_ln()
        for he in ln.findall(qn('a:headEnd')):
            ln.remove(he)
        he = etree.SubElement(ln, qn('a:headEnd'))
        he.set('type', 'triangle')
        he.set('w', 'med')
        he.set('len', 'med')
    return conn


def vline(x, y1, y2, **kw):    line(x, y1, x, y2, **kw)
def hline(x1, x2, y, **kw):    line(x1, y, x2, y, **kw)
def varrow(x, y1, y2, **kw):   line(x, y1, x, y2, arrow=True, **kw)


def fan_out(fx, fy_bot, hub_y, cxs, box_top):
    """Bifurcación: trunk → horizontal → N ramas con flecha."""
    vline(fx, fy_bot, hub_y)
    hline(cxs[0], cxs[-1], hub_y)
    for cx in cxs:
        varrow(cx, hub_y, box_top)


def fan_in(cxs, box_bot, hub_y, tx, to_top):
    """Convergencia: N ramas → horizontal → trunk con flecha."""
    for cx in cxs:
        vline(cx, box_bot, hub_y)
    hline(cxs[0], cxs[-1], hub_y)
    varrow(tx, hub_y, to_top)


def simple(x, y1, y2):
    """Flecha simple vertical."""
    varrow(x, y1, y2)


def label(x, y, w, h, text, fs=9.5, bold=False, italic=True, color=None, align=PP_ALIGN.CENTER):
    color = color or C_DGRAY
    tb = slide.shapes.add_textbox(Cm(x), Cm(y), Cm(w), Cm(h))
    tf = tb.text_frame
    tf.word_wrap = False
    p = tf.paragraphs[0]
    p.alignment = align
    r = p.add_run()
    r.text = text
    r.font.name = 'Calibri'
    r.font.size = Pt(fs)
    r.font.bold = bold
    r.font.italic = italic
    r.font.color.rgb = color


# ══════════════════════════════════════════════════════════════════════════════
# TÍTULO
# ══════════════════════════════════════════════════════════════════════════════
label(0.4, 0.08, 27.2, 0.60,
      "Pipeline de Inferencia de Elección Modal",
      fs=16, bold=True, italic=False, color=C_BLACK)

# ══════════════════════════════════════════════════════════════════════════════
# B1 — Entrada (cian)
# ══════════════════════════════════════════════════════════════════════════════
box(CX - MAIN_W/2, Y_B1, MAIN_W, MAIN_H,
    C_CYAN,
    "POST   /api/lpmc/predict   |   /compare",
    "origen  ·  destino  ·  perfil sociodemográfico",
    ts=12.5, ss=9.5)

# ── trunk de bifurcación ──────────────────────────────────────────────────────
vline(CX, B1_BOT, Y_HUB1)

# ══════════════════════════════════════════════════════════════════════════════
# FILA PARALELA 1 — asyncio.gather (OSRM ×3 + OTP)
# ══════════════════════════════════════════════════════════════════════════════
PAR_W   = 5.0
PAR_GAP = 0.47
TOTAL1  = 4 * PAR_W + 3 * PAR_GAP        # = 21.41 cm
PAR1_X0 = CX - TOTAL1 / 2               # = 3.295

# Recuadro de agrupación asyncio.gather
GROUP_PAD = 0.22
dashed_box(
    PAR1_X0 - GROUP_PAD,
    Y_PAR1  - GROUP_PAD,
    TOTAL1  + 2 * GROUP_PAD,
    H_PAR1  + 2 * GROUP_PAD,
    C_ORANGE,
    "asyncio.gather  —  concurrente"
)

PAR1_DATA = [
    ("OSRM car",  ":5000", C_ORANGE),
    ("OSRM bike", ":5001", C_ORANGE),
    ("OSRM foot", ":5002", C_ORANGE),
    ("OTP",       ":8080", C_ORANGE),
]
cxs1 = []
for i, (lbl, port, col) in enumerate(PAR1_DATA):
    bx = PAR1_X0 + i * (PAR_W + PAR_GAP)
    cx = bx + PAR_W / 2
    cxs1.append(cx)
    box(bx, Y_PAR1, PAR_W, H_PAR1, col, lbl, port, ts=12.5, ss=11)

# Fan-out + fan-in
fan_out(CX, B1_BOT, Y_HUB1, cxs1, Y_PAR1)
fan_in(cxs1, PAR1_BOT, Y_HUB2, CX, Y_B3)

# ══════════════════════════════════════════════════════════════════════════════
# B3 — build_route_features
# ══════════════════════════════════════════════════════════════════════════════
box(CX - MAIN_W/2, Y_B3, MAIN_W, MAIN_H,
    C_BGRAY,
    "build_route_features",
    "duración (s) → horas  ·  legs OTP → dur_pt_*  ·  transbordos",
    ts=12.5, ss=9.5)

simple(CX, B3_BOT, Y_B4)

# ══════════════════════════════════════════════════════════════════════════════
# B4 — build_feature_frame
# ══════════════════════════════════════════════════════════════════════════════
box(CX - MAIN_W/2, Y_B4, MAIN_W, MAIN_H,
    C_BGRAY,
    "build_feature_frame",
    "+ variables sociodemográficas  ·  one-hot purpose / fueltype",
    ts=12.5, ss=9.5)

simple(CX, B4_BOT, Y_B5)

# ══════════════════════════════════════════════════════════════════════════════
# B5 — StandardScaler
# ══════════════════════════════════════════════════════════════════════════════
box(CX - MAIN_W/2, Y_B5, MAIN_W, MAIN_H,
    C_DKGR,
    "StandardScaler.transform",
    "columnas continuas  ·  μ y σ ajustados sobre el conjunto train",
    ts=12.5, ss=9.5)

# trunk fan-out modelos
vline(CX, B5_BOT, Y_HUB3)

# ══════════════════════════════════════════════════════════════════════════════
# FILA PARALELA 2 — Modelos
# ══════════════════════════════════════════════════════════════════════════════
MOD_W   = 5.8
MOD_GAP = 0.9
TOTAL2  = 3 * MOD_W + 2 * MOD_GAP       # = 19.2 cm
MOD_X0  = CX - TOTAL2 / 2              # = 4.4

MOD_DATA = [
    ("XGBoost",          C_DKBL),
    ("Random Forest",    C_GREEN),
    ("DNN  —  PyTorch",  C_RED),
]
cxs2 = []
for i, (lbl, col) in enumerate(MOD_DATA):
    bx = MOD_X0 + i * (MOD_W + MOD_GAP)
    cx = bx + MOD_W / 2
    cxs2.append(cx)
    box(bx, Y_PAR2, MOD_W, H_PAR2, col, lbl, ts=14)

fan_out(CX, B5_BOT, Y_HUB3, cxs2, Y_PAR2)
fan_in(cxs2, MOD_BOT, Y_HUB4, CX, Y_B7)

# ══════════════════════════════════════════════════════════════════════════════
# B7 — Respuesta JSON (cian)
# ══════════════════════════════════════════════════════════════════════════════
box(CX - MAIN_W/2, Y_B7, MAIN_W, MAIN_H,
    C_CYAN,
    "Respuesta JSON",
    "predicted_mode  ·  confidence  ·  probabilities [walk, cycle, pt, drive]",
    ts=12.5, ss=9.5)

# ── Guardar ───────────────────────────────────────────────────────────────────
prs.save("pipeline_inferencia.pptx")
print("Guardado: pipeline_inferencia.pptx")
